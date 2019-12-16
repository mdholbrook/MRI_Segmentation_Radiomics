import os
import json
from glob2 import glob
import nibabel as nib
import numpy as np
from datetime import datetime
from skimage.morphology import binary_dilation, ball
import tensorflow as tf
from skimage.transform import resize
from time import time
import pptx
from pptx import util
from pylab import rcParams
import seaborn as sns
from PIL import Image
from skimage import color


# Set up plotting properties
sns.set(style='ticks', palette='Spectral', font_scale=1.5)
rcParams['figure.figsize'] = 6, 4


def load_json(base_path):
    """
    Load processing log
    Args:
        base_path (str): base directory

    Returns:

    """

    # Get file path
    json_file = os.path.join(base_path, 'processing_log.json')

    # Read file
    with open(json_file, 'r') as f:
        df = json.load(f)

    return df


def animal_files(base_path, animal_id, study_date):
    """
    Get image and segmentation files for each animal
    Args:
        base_path (str): base path
        animal_id (str): animal ID
        study_date (str): study date

    Returns:
        (str): image path
        (str): mask path
    """

    # Get animal path
    an_path = os.path.join(base_path, animal_id)

    # Get study date
    st_path = os.path.join(an_path, study_date + '*')
    st_path = glob(st_path)[0]

    # Get T2 image
    t2 = glob(os.path.join(st_path, 'T2_cor.nii*'))[0]

    # Get mask image
    mask = glob(os.path.join(st_path, 'tumor_seg*'))[0]

    return t2, mask


def dilate_mask(mask, dilator, diff):

    # Create dilation kernel
    selem = ball(dilate_kernel_size)

    # Account for anisotropic resolution 100 um x 100 um x 300 um
    sz = selem.shape
    selem = resize(selem, (sz[0], sz[1], sz[2]//3))
    selem = selem > 0
    selem = selem.astype(np.uint8)

    # Perform dilation
    mask_dia = binary_dilation(mask, selem)
    t = time()
    mask_dia = dilator.compute(mask)

    if diff:
        # Only return the tumor boundary, otherwise use tumor bed
        mask_dia = np.logical_xor(mask, mask_dia)

    print('\t\tTime to dilate: %0.3f seconds' % (time() - t))

    return mask_dia


def display_segmentations(t2, y_pred, save_path, sname='segs.png'):
    """
    Saves segmentations as an overlayed montage.
    Args:
        t2 (3D numpy array): T2 image
        y_pred (3D numpy array): Binary segmentation
        save_path (str): directory in which to save images

    Returns:

    """

    # Make mask see-through
    y_mask = y_pred.squeeze()

    # Set up slices to plot
    ims = 4
    slices = range(0, y_pred.shape[2], ims)
    rows = 3
    cols = len(slices) // rows
    resh = (y_pred.shape[0] * rows, y_pred.shape[1] * cols)

    t2_im = np.zeros(shape=(resh))
    y_mask_im = np.zeros_like(t2_im)
    r = [0, y_pred.shape[0]]
    c = [0, y_pred.shape[0]]
    n = 0
    for i in range(rows):
        c = [0, y_pred.shape[0]]

        for ii in range(cols):
            t2_im[r[0]:r[1], c[0]:c[1]] = t2[:, :, slices[n]].T
            y_mask_im[r[0]:r[1], c[0]:c[1]] = y_mask[:, :, slices[n]].T
            n += 1
            c = [i + y_pred.shape[0] for i in c]

        r = [i + y_pred.shape[0] for i in r]

    # Mask out background
    y_mask_im = np.ma.masked_where(y_mask_im.astype(bool) == 0, y_mask_im.astype(bool))

    # Save images
    # https://stackoverflow.com/questions/9193603/applying-a-coloured-overlay-to-an-image-in-either-pil-or-imagemagik
    # Convert images to RGB
    t2_im_rep = np.dstack((t2_im, t2_im, t2_im))
    y_mask_im_rep = np.zeros_like(t2_im_rep)
    y_mask_im_rep[:, :, 1] = y_mask_im

    # Convert the input image and color mask to Hue Saturation Value (HSV)
    # colorspace
    y_mask_im_hsv = color.rgb2hsv(y_mask_im_rep)
    t2_im_hsv = color.rgb2hsv(t2_im_rep)

    # Replace the hue and saturation of the original image
    # with that of the color mask
    alpha = 0.6
    t2_im_hsv[:, :, 0] = y_mask_im_hsv[:, :, 0]
    t2_im_hsv[:, :, 1] = y_mask_im_hsv[:, :, 1] * alpha

    # Convert bach to RGB
    im_masked = color.hsv2rgb(t2_im_hsv)

    # Convert image to 8-bit
    im_masked -= im_masked.min()
    im_masked /= im_masked.max() * 0.8
    im_masked *= 255
    im_masked = im_masked.astype(np.uint8)

    # Save as image using PIL
    im = Image.fromarray(im_masked)

    # Resize image for easier viewing
    im_size = (im.width//2, im.height//2)
    im = im.resize(im_size, Image.BICUBIC)
    sname = os.path.join(save_path, sname)
    im.save(sname, 'png')

    return sname


class mask_dilation_class:
    def __init__(self, kernel_size, gpu='GPU:0'):
        self.mask_in = []
        self.dia_graph = []
        self.gpu = gpu

        selem = ball(kernel_size)

        # Account for anisotropic resolution 100 um x 100 um x 300 um
        sz = selem.shape
        selem = resize(selem, (sz[0], sz[1], sz[2] // 3))
        selem = selem > 0
        selem = selem.astype(np.uint8)
        selem = selem[:, :, :, np.newaxis, np.newaxis]

        self.kernel = selem

        self.model = self.graph_keras(self.kernel)

    def graph_keras(self, kernel):

        with tf.device(self.gpu):
            # Inputs
            x_input = tf.keras.layers.Input(shape=(None, None, None, 1), name='x_input')

            # Filter input image
            lp_input = tf.keras.layers.Conv3D(filters=1, padding='SAME', kernel_size=(kernel.shape[0],
                                                                                      kernel.shape[1],
                                                                                      kernel.shape[2]),
                                        use_bias=False, trainable=False, name='dilate')(x_input)

            # Assemble model
            model = tf.keras.models.Model(inputs=x_input, outputs=lp_input)

            # Set weights for non-trainable layer
            model.layers[1].set_weights([kernel])

            model.compile(loss='mean_squared_error', optimizer=tf.keras.optimizers.Adam())

        return model

    def compute(self, mask):

        mask = mask[np.newaxis, :, :, :, np.newaxis]
        dia = self.model.predict(mask)
        dia = np.squeeze(dia)

        return dia.astype('bool')


if __name__ == '__main__':

    base_path = '/media/matt/Seagate Expansion Drive/b7TData_19/b7TData/Results'
    dilate_kernel_size = 25

    # Get the processing data
    log = load_json(base_path)

    # Compute number of animals
    n_animals = len(log)
    animals = log.keys()

    # Set up tensorflow for dilation
    dilator = mask_dilation_class(dilate_kernel_size)

    # Set up PowerPoint
    ppt_file = os.path.join(base_path, 'sarcoma_segmentations.pptx')
    title_slide = 0
    subtitle_slide = 2
    title_and_content = 5
    if not os.path.exists(ppt_file):
        prs = pptx.Presentation()

        title_slide_layout = prs.slide_layouts[title_slide]
        slide = prs.slides.add_slide(title_slide_layout)

        title = slide.shapes.title
        title.text = 'Sarcoma Segmentations'

        subtitle = slide.placeholders[1]
        subtitle.text = 'Created %s' % datetime.strftime(datetime.now(), '%B %d, %Y')

    else:
        prs = pptx.Presentation(ppt_file)

    for animal in animals:

        print('Processing %s' % animal)

        # Get study date information
        study_dates = log[animal]['StudyDate']
        n_study_dates = len(study_dates)

        # Sort dates chronologically
        study_dates = np.sort(study_dates)

        # Create animal subheading
        subsection = prs.slide_layouts[subtitle_slide]
        slide = prs.slides.add_slide(subsection)

        title = slide.shapes.title
        title.text = animal

        # Get RT status
        if len(study_dates) == 2:
            RT_flags = ['PreRT', 'PostRT']
        elif len(study_dates) == 1:
            RT_flags = ['', '']

        for j, study_date in enumerate(study_dates):

            print('\tStudy date %s' % study_date)

            # Get file names
            t2_file, mask_file = animal_files(base_path, animal, study_date)

            # Load images
            t2 = nib.load(t2_file).get_data().astype('int16')
            mask = nib.load(mask_file).get_data().astype(bool)

            # Get post/pre RT
            RT_flag = RT_flags[j]

            # Create dilated segmentations
            for i in range(1):

                if i == 0:

                    mask_dia = mask.copy()
                    section_title = '%s %s - tumor mask' % (study_date, RT_flag)

                elif i == 1:

                    mask_dia = dilate_mask(mask, dilator, diff=False)
                    section_title = '%s %s - tumor bed' % (study_date, RT_flag)

                else:

                    mask_dia = dilate_mask(mask, dilator, diff=True)
                    section_title = '%s %s - tumor margin' % (study_date, RT_flag)

                # Create montage images
                mnames = display_segmentations(t2, mask_dia, save_path=os.getcwd(), sname='segs%d.png' % i)

                # Update PowerPoint
                content_slide = prs.slide_layouts[title_and_content]
                slide = prs.slides.add_slide(content_slide)

                title = slide.shapes.title
                title.text = section_title

                left = util.Inches(0.2)
                top = util.Inches(1.5)
                width = util.Inches(9.5)
                height = util.Inches(5.7)
                pic = slide.shapes.add_picture(mnames, left, top, width, height)

                # Delete montage images
                os.remove(mnames)

        prs.save(ppt_file)


