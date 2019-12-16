"""
A collection of functions used to analyze radiomic features. This code appends data from RT, primary tumor
recurrence, and survival times. The example code is at the end of this file.
"""

import os
import sys
file = '/home/matt/Documents/SegSarcoma'
sys.path.append(file)
import pptx
import pickle
import json
from datetime import datetime, timedelta
import datetime as dt
from time import time
from subprocess import call
from lifelines import KaplanMeierFitter
import numpy as np
import pandas as pd
from matplotlib import pyplot as plt
from matplotlib import rcParams
import seaborn as sns
from glob2 import glob
from scipy import stats
from Crawler.crawler_radiomics import load_study_data
from Radiomics.radiomic_functions import *
from PIL import Image, ImageDraw, ImageFont

# Set up plotting properties
sns.set()
sns.set_style('whitegrid')
sns.set_context("paper")

PATH_TO_MRMR = '/home/matt/Downloads/mrmr_c_src/mrmr'


def concate_contrasts_small(df):
    """
    Takes a dataframe containing radiomcis data for all animals and contrasta.
    Returns a dataframe with contrasts of the same animal concatenated together.
    Args:
        df:

    Returns:

    """

    # Turn multiple contrast images into a single vector of data
    keys = df.keys()
    cat_cols = [i + '_T1' for i in keys] + \
               [i + '_T1C' for i in keys] + \
               [i + '_T2' for i in keys]
    key = [[i + '_T1' for i in keys],
           [i + '_T1C' for i in keys],
           [i + '_T2' for i in keys]]
    df_cat = pd.DataFrame(columns=cat_cols, index=range(1))

    ind = 0
    for i in range(3):

        tmp = df.iloc[i]

        # Append data into larger array
        df_cat[key[i]] = [n for n in tmp.get_values().ravel()]

    return df_cat


def strip_excluded(rfiles, exclude):
    """
    Remove animals which either did not have two scans or died shortly after surgery.
    Args:
        rfiles (list): list of all animals
        exclude (list): list of animals to exclude

    Returns:
        (list): filtered list
    """

    out_rfiles = rfiles

    for file in rfiles:

        path, _ = os.path.split(file)
        path, _ = os.path.split(path)
        _, animal = os.path.split(path)

        if animal in exclude:
            out_rfiles.remove(file)

    return out_rfiles


def keep_animals(rfiles, keep):
    """
    Only keep specific animals.
    Args:
        rfiles (list): list of all animals
        keep (list): list of animals to keep

    Returns:
        (list): filtered list
    """

    out_rfiles = []

    for file in rfiles:

        path, _ = os.path.split(file)
        path, _ = os.path.split(path)
        _, animal = os.path.split(path)

        if animal in keep:
            out_rfiles.append(file)

    return out_rfiles


def select_region(rfiles, region):
    """
    Select features based on the region they were calculated in.
    Args:
        rfiles (list): list of all animals
        region (str): 'tumor', 'bed', or 'edge'

    Returns:

    """

    out_rfiles = rfiles

    for _ in range(2):

        for file in rfiles:

            if region == 'tumor':

                if '_bed' in file:
                    out_rfiles.remove(file)
                elif 'edge' in file:
                    out_rfiles.remove(file)

            elif region == 'bed':

                if '_bed' not in file:
                    out_rfiles.remove(file)

            elif region == 'edge':

                if '_edge' not in file:
                    out_rfiles.remove(file)

        rfiles = out_rfiles

    return out_rfiles


def load_radiomics(radiomics_paths, exclude, region, keep, group):
    """
    Load radimic data (output of PyRadiomics)
    Args:
        radiomics_paths (str): path to radiomic files
        exclude (str): amimals to exclude
        region (str): which region to select
        keep (str): animals to keep
        group (str): pre/post RT and control/treatment

    Returns:
        pandas dataframe: radiomic data
    """

    # Load radiomics file data
    with open(radiomics_paths, 'r') as f:
        rfiles = f.readlines()

    # Remove line endings '\n'
    rfiles = [i.strip() for i in rfiles]

    # Sort radiomics files
    rfiles = strip_excluded(rfiles, exclude)
    rfiles = select_region(rfiles, region)

    # Only load specific animals
    rfiles = keep_animals(rfiles, keep)

    # Load data and combine contrasts
    df = load_and_concate(rfiles)

    df['Group'] = group
    # df['Region'] = region

    df.index = range(df.shape[0])

    return df


def load_study_logs(summary_file, log_file):
    """
    Load study logs which contain relevant dates
    Args:
        summary_file:
        log_file:

    Returns:

    """

    # Load summary file
    summary = load_study_data(summary_file)

    # Convert animal IDs to match the log file
    summary['Kirsch lab iD'] = ['K' + str(int(i)) for i in summary['Kirsch lab iD']]

    if not log_file == '':
        # Load log file
        with open(log_file) as f:
            log = json.load(f)
    else:
        log = 0

    # Amputation date
    treat_date = summary['Date of Antibody treatment   (0.2mg/mouse Ip)  1st']

    # Compute time until today
    summary['SinceAmputation'] = (datetime.now() - treat_date)

    return summary, log


def capitalize_method(key):
    """
    Capitalizes an list of strings
    Args:
        key (str): a feature name to capitalize

    Returns:
        (str): capitalized name
    """

    if str in key:

        key = key.replace(str, str.upper())

    return key


def load_and_concate(rfiles):
    """
    Loads and concatenates radiomic features from different MR contrasts
    Args:
        rfiles (list): list of radiomic files (one for each scan)

    Returns:
        pandas dataframe: radiomic features
    """

    animal_id = []
    for i, file in enumerate(rfiles):

        df = pd.DataFrame.from_csv(file)

        if i == 0:
            cat_df = concate_contrasts_small(df)

        else:
            df = concate_contrasts_small(df)
            cat_df = pd.concat([cat_df, df])

        path, _ = os.path.split(file)
        path, _ = os.path.split(path)
        _, animal = os.path.split(path)

        if animal not in animal_id:
            animal_id.append(animal)

        else:
            print(animal)

    # Filter out non-feature names
    features = cat_df.keys()
    feature_names = list(
        sorted(filter(lambda k: k.startswith("original_"), features)))
    cat_df = cat_df[feature_names]

     # Add animal name
    cat_df['ID'] = animal_id

    # Update df index
    cat_df.index = range(cat_df.shape[0])

    return cat_df


def load_recurrence_log(recurrence_file, threshold):
    """
    Load information about recurrence
    Args:
        recurrence_file (str): file containing primary tumor recurrence information.
        threshold (int): a limit on how long a mouse must survive post-surgury to be included

    Returns:
        pandas dataframe: animal and recurrence information
    """

    # Read in log file
    dat = pd.read_excel(recurrence_file)

    # Add 'K' to each animal name
    dat['Animal'] = ['K' + str(n) for n in dat['Animal']]

    # Filter out animals without amputation
    dat = dat[dat['AmputationDate'].notnull()]

    # Fill in death dates with today
    dat['DeathDate'][dat['DeathDate'].isnull()] = datetime.now()

    # Find lifespan post amputation
    dat['AmputationDate'] = [i.date() for i in dat['AmputationDate']]
    dat['DeathDate'] = [i.date() for i in dat['DeathDate']]
    dat['Recurrence'] = [i.date() for i in dat['Recurrence']]
    lifespan = dat['DeathDate'] - dat['AmputationDate']
    dat['since_amp'] = [int(life.days) for life in lifespan]

    # Filter out short lifespans without recurrence
    thresh = timedelta(threshold)
    recurr = ~pd.isna(dat['Recurrence'])
    inds = np.logical_or(lifespan > thresh, recurr)
    dat = dat[inds]
    # dat = dat[lifespan > thresh]

    # Get recurrence
    dat['bool_recur'] = dat['Recurrence'].notnull().tolist()

    # Reindex
    dat = dat.reset_index(drop=True)

    return dat


def sort_study_data(rec_df, sum_df, exclude):

    # Create a new dictionary to sort recurrence data
    df = {'animalID': [], 'group': [], 'recurrence': [], 'rec_days': []}

    # Populate df with data from both logs
    n = 0
    for i, animal in enumerate(sum_df['Kirsch lab iD']):

        if any(animal == rec_df['Animal']) and animal not in exclude:

            tmp = rec_df[rec_df['Animal'].str.match(animal)]

            df['animalID'].append(animal)
            df['group'].append(sum_df['Group'][i])
            df['recurrence'].append(tmp['bool_recur'].tolist()[0])
            df['rec_days'].append(int(tmp['since_amp']))

    # Make a Pandas df
    df = pd.DataFrame.from_dict(df)

    # Print general stats
    nrec = sum(df['recurrence'])
    norec = sum(df['recurrence']==False)

    nrecpd1 = sum([df['group'][i] == 'PD1' and df['recurrence'][i] for i in range(len(df))])
    nreccon = sum([df['group'][i] == 'Control' and df['recurrence'][i] for i in range(len(df))])

    norecpd1 = sum([df['group'][i] == 'PD1' and not df['recurrence'][i] for i in range(len(df))])
    noreccon = sum([df['group'][i] == 'Control' and not df['recurrence'][i] for i in range(len(df))])

    print('Number of recurrent tumors:     %d' % nrec)
    print('\tControl:    \t%d' % nreccon)
    print('\tPD1:        \t%d' % noreccon)
    print('Number of non-recurrent tumors: %d' % norec)
    print('\tControl:    \t%d' % noreccon)
    print('\tPD1:        \t%d' % norecpd1)

    return df


def append_rec(df_rec, df_tumor):
    """
    Append radiomics, RT, survival, and recurrence data together
    Args:
        df_rec (pandas dataframe): recurrence
        df_tumor (pandas dataframe): radiomics

    Returns:
        appended pandas dataframe
    """

    # Get the list of animals
    animals = df_rec['animalID'].tolist()
    keys = df_rec.keys().tolist()
    keys.remove('animalID')
    keys.remove('group')

    # Extend the output dataframe
    ltumor = len(df_tumor['ID'])
    df_out = df_tumor.copy()
    for key in keys:
        df_out[key] = pd.Series(np.zeros(ltumor), index=df_out.index)

    for animal in animals:

        rec = df_rec[df_rec['animalID'] == animal]

        for key in keys:
            df_out.loc[df_out['ID'] == animal, key] = rec[key].tolist()

    return df_out


def examine_recurrence(df_tumor, spath, group='Pre'):
    """
    Run classifiers to predict recurrence
    Args:
        df_tumor (pandas dataframe): radiomic data
        spath (str): output path
        group (str): pre/post RT

    Returns:

    """

    # Ensure spath exists
    if not os.path.exists(spath):
        os.makedirs(spath)

    # Separate control and PD1
    df_cnt = df_tumor[df_tumor['Group'] == group + 'Cnt']
    df_pd1 = df_tumor[df_tumor['Group'] == group + 'PD1']

    # Separate by larger pre/post
    df_tumor = df_tumor.loc[(df_tumor['Group'] == group + 'Cnt') | (df_tumor['Group'] == group + 'PD1'), :]

    # Get numeric data
    keys = df_tumor.keys().tolist()
    rm_keys = ['ID', 'Group', 'rec_days']
    [keys.remove(i) for i in rm_keys]
    lab_keys = 'recurrence'
    data_keys = keys
    data_keys.remove(lab_keys)

    # Scale measurements
    rad = df_tumor[data_keys].copy()
    # rad = scale_features(rad)
    label = df_tumor[lab_keys]
    ids = df_tumor['ID']
    # bar_plot(rad_data)

    # tmp = rad.loc[label, 'original_firstorder_10Percentile_T1']
    # rad.loc[label, 'original_firstorder_10Percentile_T1'] = tmp.copy()

    # Plot metrics
    plot_radiomics(rad, label, ids, spath)

    # SVM classification
    svm_classifier_loc(scale_features(rad), label, spath)
    # tspath = os.path.join(spath, 'cnt')
    # svm_classifier_loc(scale_features(df_cnt[data_keys]), df_cnt[lab_keys], tspath)
    # tspath = os.path.join(spath, 'pd1')
    # svm_classifier_loc(scale_features(df_pd1[data_keys]), df_pd1[lab_keys], tspath)

    # NN classification
    nn_classifier_loc(scale_features(rad), label, spath)
    # tspath = os.path.join(spath, 'cnt')
    # nn_classifier_loc(scale_features(df_cnt[data_keys]), df_cnt[lab_keys], tspath)
    # tspath = os.path.join(spath, 'pd1')
    # nn_classifier_loc(scale_features(df_pd1[data_keys]), df_pd1[lab_keys], tspath)

    # sname = os.path.join(spath, 'full.csv')
    # pd.concat((scale_features(rad), label), axis=1).to_csv(sname, index=False)
    # sname = os.path.join(spath, 'cnt.csv')
    # pd.concat((scale_features(df_cnt[data_keys]), df_cnt[lab_keys]), axis=1).to_csv(sname, index=False)
    # sname = os.path.join(spath, 'pd1.csv')
    # pd.concat((scale_features(df_pd1[data_keys]), df_pd1[lab_keys]), axis=1).to_csv(sname, index=False)


def svm_classifier_loc(rad, label, spath):
    """
    Attempts at SVM classification using the full dataset and PCA-based reduced datasets.
    Args:
        rad (pandas dataframe): radiomic features
        label (pandas dataframe): classification label
        spath (str): output directory

    Returns:

    """

    sns.set(font_scale=1.1)

    # Ensure spath exists
    if not os.path.exists(spath):
        os.makedirs(spath)

    old_stdout = sys.stdout
    log_file_path = os.path.join(spath, 'SVM_results.txt')
    log_file = open(log_file_path, 'w')
    sys.stdout = log_file

    print('Training samples: %d, test samples: %d' % (int(rad.shape[0]*0.75), int(rad.shape[0]*0.25)))
    print('Number of features %d' % rad.shape[1])

    # Get index
    inds_1 = label == True
    inds_2 = label == False

    # Create label vector
    label = label.get_values().astype(np.float64)
    rad = rad.get_values().astype(np.float64)

    # Plot euclidean distances - multidimensional scaling
    fig = plot_euclidean_distances(rad, inds_1, inds_2)
    fig.savefig(os.path.join(spath, 'SVM_EuclDist_SVM.png'), dpi=300)
    plt.close(fig)

    # SVM without reduction
    print('Trying SVM classifier without dimensionality reduction:')
    print('-' * 50)
    # fig, clf = svm_classifier(rad.get_values(), inds_1, inds_2)
    fig = svm_cross(rad, label)
    fig.savefig(os.path.join(spath, 'SVM_NoPCA_roc.svg'))
    fig.savefig(os.path.join(spath, 'SVM_NoPCA_roc.png'), dpi=200)
    plt.close(fig)

    with open(os.path.join(spath, 'NoReduction_SVM.model'), 'wb') as f:
        pickle.dump(clf, f)

    # PCA analysis
    # a look at 2 components and SVM
    print('\n')
    print('Trying SVM classifier with 2 principal components:')
    print('-' * 50)
    rad_pca = pca_reduction(rad, npcomps=2)
    fig = svm_cross(rad_pca, label)
    fig.savefig(os.path.join(spath, 'SVM_PCA.svg'))
    plt.close(fig)
    fig = svm_cross(rad_pca, label)
    fig.savefig(os.path.join(spath, 'SVM_PCA2_roc.svg'))
    fig.savefig(os.path.join(spath, 'SVM_PCA2_roc.png'), dpi=300)
    plt.close(fig)
    with open(os.path.join(spath, 'PCA2_SVM.model'), 'wb') as f:
        pickle.dump(clf, f)

    # a look at variance
    print('\n')
    npcomps, fig = pca_analyis(rad)
    print('Trying SVM classifier with %d principal components found from data variance:' % npcomps)
    print('-' * 50)
    rad_pca = pca_reduction(rad, npcomps)
    fig.savefig(os.path.join(spath, 'Cumulative_var.svg'))
    # fig.savefig(os.path.join(spath, 'Cumulative_var.png'), dpi=300)
    plt.close(fig)

    fig = svm_cross(rad_pca, label)
    fig = svm_cross(rad_pca, label)
    # fig.savefig(os.path.join(spath, 'SVM_PCA_roc_var.svg'))
    # fig.savefig(os.path.join(spath, 'SVM_PCA_roc_var.png'), dpi=300)
    plt.close(fig)
    with open(os.path.join(spath, 'SVM_PCA95.model'), 'wb') as f:
        pickle.dump(clf, f)

    # Close log file
    sys.stdout = old_stdout
    log_file.close()


def nn_classifier_loc(rad, label, spath):
    """
    Attempts at NN classification using the full dataset and PCA-based reduced datasets.
    Args:
        rad (pandas dataframe): radiomic features
        label (pandas dataframe): classification label
        spath (str): output directory

    Returns:

    """

    sns.set(font_scale=1.1)

    # Ensure spath exists
    if not os.path.exists(spath):
        os.makedirs(spath)

    old_stdout = sys.stdout
    log_file_path = os.path.join(spath, 'NN_results.txt')
    log_file = open(log_file_path, 'w')
    sys.stdout = log_file

    print('Training samples: %d, test samples: %d' % (int(rad.shape[0]*0.75), int(rad.shape[0]*0.25)))
    print('Number of features %d' % rad.shape[1])

    # Get index
    inds_1 = label == True
    inds_2 = label == False

    # Create label vector
    label = label.get_values().astype(np.float64)
    rad = rad.get_values().astype(np.float64)

    # SVM without reduction
    print('Trying NN classifier without dimensionality reduction:')
    print('-' * 50)
    # fig, clf = svm_classifier(rad.get_values(), inds_1, inds_2)
    fig = neural_network_cross(rad, label)
    fig.savefig(os.path.join(spath, 'NN_NoPCA_roc.svg'))
    fig.savefig(os.path.join(spath, 'NN_NoPCA_roc.png'), dpi=200)
    plt.close(fig)

    # PCA analysis
    # a look at 2 components and SVM
    print('\n')
    print('Trying NN classifier with 2 principal components:')
    print('-' * 50)
    rad_pca = pca_reduction(rad, npcomps=2)
    # fig, clf = svm_classifier(rad_pca, inds_1, inds_2)
    # fig.savefig(os.path.join(spath, 'PCA_NN.png'), dpi=300)
    # plt.close(fig)
    fig = neural_network_cross(rad_pca, label)
    # fig.savefig(os.path.join(spath, 'NN_PCA2_roc.svg'))
    # fig.savefig(os.path.join(spath, 'NN_PCA2_roc.png'), dpi=300)
    plt.close(fig)
    # with open(os.path.join(spath, 'PCA2_SVM.model'), 'wb') as f:
    #     pickle.dump(clf, f)

    # a look at variance
    print('\n')
    npcomps, fig = pca_analyis(rad)
    print('Trying NN classifier with %d principal components found from data variance:' % npcomps)
    print('-' * 50)
    rad_pca = pca_reduction(rad, npcomps)
    # fig.savefig(os.path.join(spath, 'Cumulative_var.png'), dpi=300)
    plt.close(fig)

    # fig, clf = svm_classifier(rad_pca, inds_1, inds_2)
    # fig = neural_network_cross(rad_pca, label)
    # fig.savefig(os.path.join(spath, 'NN_PCA_var_roc.svg'))
    # fig.savefig(os.path.join(spath, 'NN_PCA_var_roc.png'), dpi=300)
    # plt.close(fig)
    # with open(os.path.join(spath, 'PCA95_NN.model'), 'wb') as f:
    #     pickle.dump(clf, f)

    # Close log file
    sys.stdout = old_stdout
    log_file.close()


def mRMR(df_tumor, base_path, group='Post', area='tumor', mrmr_file=None, num_features=10, htmaps=False):
    """
    Compute mRMR feature selection. This funtion has been bundled with code which computes correlation
    maps and later calls training on classifiers for tumor recurrent.
    Args:
        df_tumor (panas dataframe):
        base_path (str): where to save mRMR output
        group (str): pre/post RT
        area (str): masked location
        mrmr_file (str): file of mRMR features to load or create
        num_features (int): number of features to use
        htmaps (bool): whether to create and save all heatmaps

    Returns:

    """

    print('\n')
    print('-'*80)
    print('Processing %s, %s' %(group, area))
    print('Number of features: %d' % num_features)

    dfn = df_tumor.loc[(df_tumor['Group'] == group + 'Cnt') | (df_tumor['Group'] == group + 'PD1'), :]
    dfn = dfn.reset_index()

    keys = df_tumor.keys().tolist()
    rm_keys = ['ID', 'Group', 'rec_days']
    [keys.remove(i) for i in rm_keys]
    lab_keys = 'recurrence'
    data_keys = keys
    data_keys.remove(lab_keys)

    # # Separate by larger pre/post
    # df_tumor = df_tumor[(df_tumor['Group'] == group + 'Cnt') | (df_tumor['Group'] == group + 'PD1')]

    # Create class in the first column
    category = dfn['recurrence']

    # Remove columns
    dfn = dfn[data_keys].copy()

    dfn['Class'] = category.astype(int)

    # Rearrange columns so class is first
    cols = dfn.columns.tolist()
    cols = [cols[-1]] + cols[:-1]
    dfn = dfn[cols]

    # Drop outliers
    dfn = dfn.loc[dfn['original_gldm_GrayLevelNonUniformity_T1'] != 5607, :]

    # Write as csv
    sname = 'tumor_%s_%s.csv' % (group, area)
    dfn.to_csv(sname, index=False, index_label=False)

    # Generate savepath
    if mrmr_file:
        if 'post' in mrmr_file.lower():
            mrmr_time = 'Post'
        elif 'pre' in mrmr_file.lower():
            mrmr_time = 'Pre'
        if 'tumor' in mrmr_file.lower():
            mode = 'tumor'
        elif 'bed' in mrmr_file.lower():
            mode = 'bed'
        else:
            mode = 'edge'
        spath = os.path.join(base_path, 'Analysis', 'mrmr_%s_%s_feat_%s_%s' % (group, area, mrmr_time, mode))
    else:
        spath = os.path.join(base_path, 'Analysis', 'mrmr_%s_%s' % (group, area))

    print('\tSaving to: ', spath)
    if not os.path.exists(spath):
        os.makedirs(spath)

    # Run mRMR
    if not mrmr_file:
        print('\tCreating a new mRMR file')

        log_file_path = os.path.join(base_path, 'Analysis', 'mRMR_results_%s_%s.txt' % (group, area))
        log_file = open(log_file_path, 'w')
        call([PATH_TO_MRMR, '-i', sname, '-t', '1', '-n', '200'],
             stdout=log_file)
        log_file.close()

        # Select out predictive features
        mrmr_file = os.path.join(base_path, 'Analysis', 'mRMR_results_%s_%s.txt' % (group, area))
    else:
        mrmr_file = os.path.join(base_path, 'Analysis', mrmr_file)
        print('\tUsing the specified mRMR file:\n\t\t%s' % mrmr_file)

    # Read mRMR features
    # res = pd.read_csv(mrmr_file, delimiter='\t', header=205, skipfooter=9, skipinitialspace=True, engine='python')
    res = pd.read_csv(mrmr_file, delimiter='\t', header=3, skipfooter=212, skipinitialspace=True, engine='python')
    cols = res.keys().tolist()
    cols = [i.strip() for i in cols]
    res.columns = cols

    features = res['Name'].tolist()
    features = [i.strip() for i in features]

    # Make heatmaps of control and pd1
    if htmaps:
        # Make heatmaps for cntr and pd1
        df_cnt = df_tumor.loc[df_tumor['Group'] == group + 'Cnt', :].copy()
        df_pd1 = df_tumor.loc[df_tumor['Group'] == group + 'PD1', :].copy()

        # Create class in the first column
        category_cnt = df_cnt['recurrence'].astype(int)
        category_pd1 = df_pd1['recurrence'].astype(int)

        # Remove columns
        df_cnt = df_cnt[features]
        df_pd1 = df_pd1[features]

        # Convert columns to float
        for key in features:
            df_cnt.loc[:, key] = pd.to_numeric(df_cnt.loc[:, key], errors='ignore')
            df_pd1.loc[:, key] = pd.to_numeric(df_pd1.loc[:, key], errors='ignore')

        # Rename to remove "original_"
        for key in df_cnt.keys():
            nkey = key.replace('original_', '')
            df_cnt = df_cnt.rename(columns={key: nkey})
            df_pd1 = df_pd1.rename(columns={key: nkey})

        # Set up colors for the animal classification
        lut = dict(zip(dfn['Class'].unique(), 'bg'))
        row_colors_cnt = category_cnt.map(lut)
        row_colors_pd1 = category_pd1.map(lut)
        row_colors_cnt.name = 'Recurrence'
        row_colors_pd1.name = 'Recurrence'

        # sns.set(font_scale=0.8)
        # g = sns.clustermap(df_cnt, figsize=(18, 18), standard_scale=1, metric='correlation', row_colors=row_colors_cnt)
        # g.savefig(os.path.join(spath, 'heatmap_animals_cnt.png'), dpi=300)
        # plt.close()
        #
        # sns.set(font_scale=0.8)
        # g = sns.clustermap(df_pd1, figsize=(18, 18), standard_scale=1, metric='correlation', row_colors=row_colors_pd1)
        # g.savefig(os.path.join(spath, 'heatmap_animals_pd1.png'), dpi=300)
        # plt.close()

        # Remove class category
        # df_cnt = df_cnt.drop('Class', axis=1)
        # df_pd1 = df_pd1.drop('Class', axis=1)

        # Compute correlation matrix
        df_cnt = df_cnt.corr()
        df_pd1 = df_pd1.corr()

        # Sort correlation matrix
        order_cnt = df_cnt.sum(axis=0).argsort()[::-1]
        order_pd1 = df_pd1.sum(axis=0).argsort()[::-1]

        # Reorder row and columns
        df_cnt = df_cnt.iloc[order_cnt, order_cnt]
        df_pd1 = df_pd1.iloc[order_cnt, order_cnt]

        # Plot correlation - control
        sns.set(font_scale=0.8)
        f, ax = plt.subplots(figsize=(18, 10))
        sns.heatmap(df_cnt, vmax=1.0, square=True, ax=ax)
        plt.yticks(rotation=0)
        ax.get_xaxis().set_ticks([])
        plt.xticks(rotation=90)
        # f.savefig(os.path.join(spath, 'heatmap_control_corr.svg'))
        f.savefig(os.path.join(spath, 'heatmap_control_corr.png'), dpi=300)
        plt.close()

        # Plot correlation - PD1
        sns.set(font_scale=0.8)
        f, ax = plt.subplots(figsize=(18, 10))
        sns.heatmap(df_pd1, vmax=1.0, square=True, ax=ax)
        plt.yticks(rotation=0)
        ax.get_xaxis().set_ticks([])
        plt.xticks(rotation=90)
        # f.savefig(os.path.join(spath, 'heatmap_pd1_corr.svg'))
        f.savefig(os.path.join(spath, 'heatmap_pd1_corr.png'), dpi=300)
        plt.close()

        # Plot correlation - difference
        sns.set(font_scale=0.8)
        f, ax = plt.subplots(figsize=(18, 10))
        sns.heatmap(np.abs(df_cnt - df_pd1), vmax=1.0, square=True, ax=ax)
        plt.yticks(rotation=0)
        ax.get_xaxis().set_ticks([])
        plt.xticks(rotation=90)
        # f.savefig(os.path.join(spath, 'heatmap_cnt_pd1_diff_corr.svg'))
        f.savefig(os.path.join(spath, 'heatmap_cnt_pd1_diff_corr.png'), dpi=300)
        plt.close()


    if htmaps:
        # Make heatmap dataframe
        htmp1 = dfn[features].copy()

        # Convert columns to float (from object)
        for key in htmp1.keys():
            htmp1.loc[:, key] = pd.to_numeric(htmp1.loc[:, key], errors='ignore')

        # Rename to remove "original_"
        for key in htmp1.keys():

            nkey = key.replace('original_', '')
            htmp1 = htmp1.rename(columns={key: nkey})

        # Set up colors for the animal classification
        lut = dict(zip(dfn['Class'].unique(), 'bg'))
        row_colors = dfn['Class'].map(lut)
        row_colors.name = 'Recurrence'

        # sns.set(font_scale=0.8)
        # g = sns.clustermap(htmp1, figsize=(18, 18), standard_scale=1, metric='correlation', row_colors=row_colors)
        # g.savefig(os.path.join(spath, 'heatmap_animals.svg'))
        # # g.savefig(os.path.join(spath, 'heatmap_animals.png'), dpi=300)
        # plt.close()

        # Compute correlation matrix - recurrence and no recurrence
        htmp_nrec = dfn.loc[dfn['Class'] == 0, features]
        htmp_rec = dfn.loc[dfn['Class'] == 1, features]

        # Convert columns to float
        for key in features:
            htmp_nrec.loc[:, key] = pd.to_numeric(htmp_nrec.loc[:, key], errors='ignore')
            htmp_rec.loc[:, key] = pd.to_numeric(htmp_rec.loc[:, key], errors='ignore')

        # Rename to remove "original_"
        for key in htmp_nrec.keys():
            nkey = key.replace('original_', '')
            htmp_nrec = htmp_nrec.rename(columns={key: nkey})
            htmp_rec = htmp_rec.rename(columns={key: nkey})

        # Compute correlation matrix
        htmp_nrec = htmp_nrec.corr()
        htmp_rec = htmp_rec.corr()

        # Sort correlation matrix
        order_nrec = htmp_nrec.sum(axis=0).argsort()[::-1]
        order_rec = htmp_rec.sum(axis=0).argsort()[::-1]

        # Reorder row and columns
        htmp_nrec = htmp_nrec.iloc[order_nrec, order_nrec]
        htmp_rec = htmp_rec.iloc[order_nrec, order_nrec]

        # Plot correlation - no recurrence
        sns.set(font_scale=0.8)
        f, ax = plt.subplots(figsize=(18, 10))
        sns.heatmap(htmp_nrec, vmax=1.0, square=True, ax=ax)
        plt.yticks(rotation=0)
        ax.get_xaxis().set_ticks([])
        plt.xticks(rotation=90)
        f.savefig(os.path.join(spath, 'heatmap_noRecur_corr.svg'))
        # f.savefig(os.path.join(spath, 'heatmap_noRecur_corr.png'), dpi=300)
        plt.close()

        # Plot correlation - recurrence
        sns.set(font_scale=0.8)
        f, ax = plt.subplots(figsize=(18, 10))
        sns.heatmap(htmp_rec, vmax=1.0, square=True, ax=ax)
        plt.yticks(rotation=0)
        ax.get_xaxis().set_ticks([])
        plt.xticks(rotation=90)
        f.savefig(os.path.join(spath, 'heatmap_Recur_corr.svg'))
        # f.savefig(os.path.join(spath, 'heatmap_Recur_corr.png'), dpi=300)
        plt.close()

        # Plot correlation - difference
        sns.set(font_scale=0.8)
        f, ax = plt.subplots(figsize=(18, 10))
        sns.heatmap(np.abs(htmp_nrec - htmp_rec), vmax=1.0, square=True, ax=ax)
        plt.yticks(rotation=0)
        ax.get_xaxis().set_ticks([])
        plt.xticks(rotation=90)
        f.savefig(os.path.join(spath, 'heatmap_Recur_diff_corr.svg'))
        # f.savefig(os.path.join(spath, 'heatmap_Recur_diff_corr.png'), dpi=300)
        plt.close()

    # Use only top 10 features
    features = features[:num_features]

    # Make heatmap dataframe - 10 features
    htmp1 = dfn[features].copy()

    # Convert columns to float (from object)
    for key in htmp1.keys():
        htmp1.loc[:, key] = pd.to_numeric(htmp1.loc[:, key], errors='ignore')

    # Rename to remove "original_"
    for key in htmp1.keys():

        nkey = key.replace('original_', '')
        htmp1 = htmp1.rename(columns={key: nkey})

    # Set up colors for the animal classification
    lut = dict(zip(dfn['Class'].unique(), 'bg'))
    row_colors = dfn['Class'].map(lut)
    row_colors.name = 'Recurrence'

    # sns.set(font_scale=0.8)
    # g = sns.clustermap(htmp1, figsize=(18, 18), standard_scale=1, metric='correlation', row_colors=row_colors)
    # g.savefig(os.path.join(spath, 'heatmap_animals.svg'))
    # # g.savefig(os.path.join(spath, 'heatmap_animals.png'), dpi=300)
    # plt.close()

    # Compute correlation matrix - recurrence and no recurrence
    htmp_nrec = dfn.loc[dfn['Class'] == 0, features]
    htmp_rec = dfn.loc[dfn['Class'] == 1, features]

    # Convert columns to float
    for key in features:
        htmp_nrec.loc[:, key] = pd.to_numeric(htmp_nrec.loc[:, key], errors='ignore')
        htmp_rec.loc[:, key] = pd.to_numeric(htmp_rec.loc[:, key], errors='ignore')

    # Rename to remove "original_"
    for key in htmp_nrec.keys():
        nkey = key.replace('original_', '')
        htmp_nrec = htmp_nrec.rename(columns={key: nkey})
        htmp_rec = htmp_rec.rename(columns={key: nkey})

    # Compute correlation matrix
    htmp_nrec = htmp_nrec.corr()
    htmp_rec = htmp_rec.corr()

    # Sort correlation matrix
    order_nrec = htmp_nrec.sum(axis=0).argsort()[::-1]
    order_rec = htmp_rec.sum(axis=0).argsort()[::-1]

    # Reorder row and columns
    htmp_nrec = htmp_nrec.iloc[order_nrec, order_nrec]
    htmp_rec = htmp_rec.iloc[order_nrec, order_nrec]

    # Plot correlation - no recurrence
    sns.set(font_scale=0.8)
    f, ax = plt.subplots(figsize=(18, 10))
    sns.heatmap(htmp_nrec, vmax=1.0, square=True, ax=ax)
    plt.yticks(rotation=0)
    ax.get_xaxis().set_ticks([])
    plt.xticks(rotation=90)
    f.savefig(os.path.join(spath, 'heatmap_noRecur_corr_10.svg'))
    # f.savefig(os.path.join(spath, 'heatmap_noRecur_corr_10.png'), dpi=300)
    plt.close()

    # Plot correlation - recurrence
    sns.set(font_scale=0.8)
    f, ax = plt.subplots(figsize=(18, 10))
    sns.heatmap(htmp_rec, vmax=1.0, square=True, ax=ax)
    plt.yticks(rotation=0)
    ax.get_xaxis().set_ticks([])
    plt.xticks(rotation=90)
    f.savefig(os.path.join(spath, 'heatmap_Recur_corr_10.svg'))
    # f.savefig(os.path.join(spath, 'heatmap_Recur_corr_10.png'), dpi=300)
    plt.close()

    # Plot correlation - difference
    sns.set(font_scale=0.8)
    f, ax = plt.subplots(figsize=(18, 10))
    sns.heatmap(np.abs(htmp_nrec - htmp_rec), vmax=1.0, square=True, ax=ax)
    plt.yticks(rotation=0)
    ax.get_xaxis().set_ticks([])
    plt.xticks(rotation=90)
    f.savefig(os.path.join(spath, 'heatmap_Recur_diff_corr_10.svg'))
    # f.savefig(os.path.join(spath, 'heatmap_Recur_diff_corr_10.png'), dpi=300)
    plt.close()

    # Create a DataFrame with mRMR features
    df_sub = df_tumor[features + rm_keys + [lab_keys]].copy()

    # Perform computations
    examine_recurrence(df_sub, spath, group=group)

    return spath


def plot_radiomics(rad, label, ids, spath):
    """
    Plot radiomic features used for recurrence classification.
    Args:
        rad (pandas dataframe): radiomic features
        label (pandas dataframe): recurrence labels
        ids (list): animal IDs
        spath (str): savepath

    Returns:

    """

    sns.set_style('whitegrid')

    # Reduce the number of features to plot to 5
    keys = rad.keys()
    rad = rad[keys[:6]]

    # Set up figure
    fig = plt.figure(figsize=(10, 6))
    ax = fig.add_axes([0.1, 0.5, 0.80, 0.47])

    # Get the list of radiomic features
    keys = rad.keys()

    # Create a dictionary (later a Dataframe) for plotting with Seaborn
    df = {'Features': [], 'Value': [], 'Recurrence': [], 'ID': []}

    # Populate dictionary, df, with features and values
    for key in keys:

        # Get feature values
        tmp = rad[key].to_numpy()

        # Normalize between -1 and 1
        tmp -= tmp.min()
        tmp /= 0.5 * tmp.max()
        tmp -= 1

        for i in range(len(tmp)):

            # Remove "original" from the plots
            nkey = key.replace('original_', '')

            # Capitalize method
            ind1 = nkey.find('_')
            method = nkey[:ind1]
            if method == 'firstorder':
                method = 'First Order'
            elif method == 'shape':
                method = method.capitalize()
            else:
                method = method.upper()

            # Get modality
            ind2 = nkey.rfind('_')
            modality = nkey[ind2 + 1:]
            if method != 'Shape':
                nkey = '%s\n%s\n%s' % (modality, method, nkey[ind1 + 1:ind2])
            else:
                nkey = '%s\n%s' % (method, nkey[ind1 + 1:ind2])


            # Append features, values, recurrence, and IDs to the dictionary
            df['Features'].append(nkey)
            df['Value'].append(tmp[i])

            if label.iloc[i]:
                df['Recurrence'].append('Yes')
            else:
                df['Recurrence'].append('No')

            df['ID'].append(ids.iloc[i])

    # Convert dictionary to Dataframe
    df = pd.DataFrame.from_dict(df)

    # Plot
    g = sns.boxplot(x='Features', y='Value', hue='Recurrence', data=df, ax=ax, palette="Set1")

    # Rotate feature names
    g.set_xticklabels(g.get_xticklabels(), rotation=60, fontsize=10)
    # g.set_yticklabels(g.get_yticklabels(), fontsize=10)
    g.set_xlabel('Features', fontsize=10)
    g.set_ylabel('Normalized Value')

    # Fix legend
    plt.legend(bbox_to_anchor=(1.01, 1), loc=2, borderaxespad=0., title='Recurrence')

    # Save figure
    fig.savefig(os.path.join(spath, 'feature_box.svg'), format='svg')
    # fig.savefig(os.path.join(spath, 'feature_box.png'), format='png', dpi=300)


def clean_key(key):
    # Remove "original" from the plots
    nkey = key.replace('original_', '')

    # Capitalize method
    ind1 = nkey.find('_')
    method = nkey[:ind1]
    if method == 'firstorder':
        method = 'First Order'
    elif method == 'shape':
        method = method.capitalize()
    else:
        method = method.upper()

    # Get modality
    ind2 = nkey.rfind('_')
    modality = nkey[ind2 + 1:]
    if method != 'Shape':
        nkey = '%s\n%s\n%s' % (modality, method, nkey[ind1 + 1:ind2])
    else:
        nkey = '%s\n%s' % (method, nkey[ind1 + 1:ind2])

    return nkey


def paired_ttests_loc(df, base_path, area):
    """
    Compute paired t-test for data
    Args:
        df (Pandas df): full radiomics data (with label columns)
        area (str): mask area

    Returns:

    """
    sns.set_style('whitegrid')

    # Set up spath
    spath = os.path.join(base_path, 'Analysis', 'ttests')
    if not os.path.exists(spath):
        os.makedirs(spath)

    # Set up output
    log_file_path = os.path.join(spath, 'ttest_%s.txt' % area)
    f = open(log_file_path, 'w')

    # Get indicies of pre/post RT
    inds_1 = df['Group'].str.contains('Pre')
    inds_2 = df['Group'].str.contains('Post')
    keys = df.keys()
    n_sets = len(df['ID'].unique())
    label = np.empty(shape=len(inds_1), dtype=str)
    non_data_labels = ['Group', 'recurrence', 'rec_days', 'ID']

    # Create a dictionary (later a Dataframe) for evaluation and plotting with Seaborn
    dfn = {'Features': [], 'P-Value': [], 'OrigFeatures': [], 'Contrast': []}


    # kkeys = [key for key in keys if key not in non_data_labels]
    # df_k = df[kkeys]
    # ob = MultiComparison(df_k, inds_1)
    # ob.allpairtest(stats.ttest_rel)

    # Populate dictionary, df, with features and values
    for key in keys:

        if key not in non_data_labels:

            # Get feature values
            tmp = df[key].to_numpy()

            p1 = tmp[inds_1]
            p2 = tmp[inds_2]

            results = stats.ttest_rel(p1, p2)

            if not np.isnan(results.pvalue):

                # Remove "original" from the plots
                nkey = key.replace('original_', '')

                # Capitalize method
                ind1 = nkey.find('_')
                method = nkey[:ind1]
                if method == 'firstorder':
                    method = 'First Order'
                elif method == 'shape':
                    method = method.capitalize()
                else:
                    method = method.upper()

                # Get modality
                ind2 = nkey.rfind('_')
                modality = nkey[ind2+1:]
                if method != 'Shape':
                    nkey = '%s\n%s\n%s' % (modality, method, nkey[ind1+1:ind2])
                else:
                    nkey = '%s\n%s' % (method, nkey[ind1+1:ind2])

                # Add contrast
                if '_T1C' in key:
                    con = 'T1C'

                elif 'T1' in key:
                    con = 'T1'

                elif 'T2' in key:
                    con = 'T2'

                else:
                    con = None

                if not nkey in dfn['Features']:
                    # Append features, values, recurrence, and IDs to the dictionary
                    dfn['Features'].append(nkey)
                    dfn['OrigFeatures'].append(key)
                    dfn['P-Value'].append(results.pvalue)
                    dfn['Contrast'].append(con)

    # Convert dictionary to Dataframe
    dfn = pd.DataFrame.from_dict(dfn)

    f.write('%d significant features found after t-tests!\n' % np.sum(dfn['P-Value'] < 0.05))

    from statsmodels.sandbox.stats.multicomp import multipletests
    mod = multipletests(dfn['P-Value'], alpha=0.05)

    # Plot corrected P-values
    fig = plt.figure(figsize=(10, 6))
    plt.plot(dfn['P-Value'], 'bo', label='Original')
    plt.plot(mod[1], 'ro', label='Corrected')
    plt.legend()
    plt.xlim([0, 400])
    plt.ylabel('P-Value')
    plt.xlabel('Features')
    # plt.savefig(os.path.join(spath, 'compare_corrected_pvalues.png'), dpi=300)
    plt.savefig(os.path.join(spath, 'compare_corrected_pvalues.svg'), dpi=300)
    plt.close(fig)

    # Remove all but the smallest p-values
    dfn['P-Value'] = mod[1]
    dfn = dfn[mod[0]]
    dfn = dfn.sort_values('P-Value', ascending=True)  # Sort by P-Value, ascending
    dfn = dfn.reset_index(drop=True)
    f.write('%d significant features found after multiple p-value correction!\n' % len(dfn))

    # p_thresh = 0.05
    # dfn = dfn.drop(dfn[dfn['P-Value'] > p_thresh].index)
    # f.write('%d significant features found!\n' % len(dfn))
    # dfn = dfn.reset_index(drop=True)

    # Get stats on which contrasts are most significant
    t1 = 0
    t1c = 0
    t2 = 0
    shape = 0
    feature_cats = {'firstorder': 0,
                    'glcm': 0,
                    'glszm':  0,
                    'glrlm': 0,
                    'glrlm': 0,
                    'ngtdm': 0,
                    'gldm': 0}

    for key in dfn['Features']:

        if 'T1C' in key:
            t1c += 1

        elif 'T1' in key:
            t1 += 1

        elif 'T2' in key:
            t2 += 1

        elif 'Shape' in key:
            shape += 1

        if 'first order' in key.lower():
            feature_cats['firstorder'] += 1
        elif 'glcm' in key.lower():
            feature_cats['glcm'] += 1
        elif 'glszm' in key.lower():
            feature_cats['glszm'] += 1
        elif 'glrlm' in key.lower():
            feature_cats['glrlm'] += 1
        elif 'ngtdm' in key.lower():
            feature_cats['ngtdm'] += 1
        elif 'gldm' in key.lower():
            feature_cats['gldm'] += 1

    counts_df = pd.DataFrame({'Contrast': ['Shape', 'T1', 'T1C', 'T2'], 'Count': [shape, t1, t1c, t2]})

    f.write('Features summary for radiomics computed from %s mask\n' % area)
    f.write('-'*40 + '\n')

    f.write('Contrast\tCount\n')
    for i, con in enumerate(counts_df['Contrast']):

        f.write('%s\t\t%d\n' % (con, counts_df['Count'].iloc[i]))

    f.write('-'*40 + '\n')
    f.write('Feature\tCount\n')
    for key in feature_cats.keys():

        f.write('%s\t\t%s\n' % (key, feature_cats[key]))



    f.write('-'*40 + '\n')
    for con, pval in zip(dfn['Features'], dfn['P-Value']):

        con = con.splitlines()
        if len(con) == 3:
            f.write('%s\n%s\n%s\nP-value: %0.5f\n\n' % (con[1], con[2], con[0], pval) )
        else:
            f.write('%s\n%s\nP-Value: %0.5f\n\n' % (con[1], con[0], pval) )


    fig = plt.figure(figsize=(10, 6))
    ax = fig.add_axes([0.11, 0.15, 0.87, 0.65])

    sns.barplot(x='Contrast', y='Count', data=counts_df, ax=ax, color='midnightblue', saturation=0.5)
    plt.ylabel('Num. features (p < 0.05)')
    # plt.show()
    fig.savefig(os.path.join(spath, 'ttest_feature_count_%s.svg' % area))

    # Plot specific categories
    cats = ['firstorder', 'shape', 'GLCM']
    mk_size = 6
    cmap = sns.color_palette()

    # Set up a df for plotting all values
    df_plot = {'Features': [], 'Value': [], 'RT': [], }

    for key, nkey in zip(dfn['OrigFeatures'], dfn['Features']):

        # Get feature values
        tmp = df[key].to_numpy()

        # Normalize between -1 and 1
        tmp -= tmp.min()
        tmp /= 0.5 * tmp.max()
        tmp -= 1

        for i in range(len(tmp)):

            # Append features, values, recurrence, and IDs to the dictionary
            df_plot['Features'].append(nkey)
            df_plot['Value'].append(tmp[i])

            if inds_1.iloc[i]:
                df_plot['RT'].append('Pre RT')
            else:
                df_plot['RT'].append('Post RT')

    df_plot = pd.DataFrame.from_dict(df_plot)

    # Set up figure
    fig = plt.figure(figsize=(20 * len(dfn['Features'])/10, 8))
    ax = fig.add_axes([0.08, 0.4, 0.9, 0.60])

    # Plot all
    g = sns.boxplot(x='Features', y='Value', hue='RT', data=df_plot, ax=ax, palette="Set1")

    # Rotate feature names
    # g.set_xticklabels(g.get_xticklabels(), rotation=35, fontsize=18)
    # ax.set_ylabel('')
    g.set_xticklabels(g.get_xticklabels(), rotation=35, fontsize=10)
    g.set_yticklabels(g.get_yticklabels(), fontsize=10)
    g.set_xlabel('Features', fontsize=10)
    ax.set_ylabel('')

    # Fix legend
    plt.legend(bbox_to_anchor=(1.01, 1), loc=2, borderaxespad=0., title='RT')

    # Save figure
    fig.savefig(os.path.join(spath, 'ttest_feature_box_%s.svg' % area), format='svg')
    # fig.savefig(os.path.join(spath, 'ttest_feature_box_%s.png' % area), format='png', dpi=300)

    # Plot top features
    nkeys = 6
    keep_keys = dfn['Features'].iloc[:nkeys]

    small_df_plot = df_plot.loc[df_plot['Features'].isin(keep_keys), :]

    fs = 10
    fig = plt.figure() #plt.figure(figsize=(10, 6))
    ax = fig.add_axes([0.08, 0.4, 0.80, 0.59])

    # Plot
    g = sns.boxplot(x='Features', y='Value', hue='RT', data=small_df_plot, ax=ax, palette="Set1")

    # Rotate feature names
    g.set_xticklabels(g.get_xticklabels(), rotation=60, fontsize=fs)
    g.set_xlabel('Features', fontsize=fs)
    g.set_ylabel('Normalized Value', fontsize=fs)

    # Fix legend
    plt.legend(bbox_to_anchor=(1.01, 1), loc=2, borderaxespad=0., title='Recurrence')

    # Save figure
    fig.savefig(os.path.join(spath, 'ttest_feature_box_select%d_%s.svg' % (nkeys, area) ), format='svg')
    # fig.savefig(os.path.join(spath, 'ttest_feature_box_select%d_%s.png' % (nkeys, area) ), format='png')

    # Plot specific features
    if 'tumor' in area:
        keep_keys = ['T2\nGLSZM\nZoneVariance',
                     'T2\nGLRLM\nRunLengthNonUniformity',
                     'Shape\nVoxelVolume',
                     'T1C\nGLSZM\nLargeAreaHighGrayLevelEmphasis',
                     'T1C\nGLDM\nDependenceNonUniformity',
                     'T1\nFirstOrder\nEnergy'
                     ]
    elif 'bed' in area:
        keep_keys = ['T2\nGLDM\nDependenceNonUniformity',
                     'T2\nFirst Order\nEnergy',
                     'Shape\nVoxelVolume',
                     'T1C\nGLRLM\nRunVariance',
                     'T1\nGLSZM\nLongRunLowGrayLevelEmphasis',
                     'T1C\nFirst Order\nRange'
                     ]
    else: # Edge
        keep_keys = ['T2\nFirst Order\nEnergy',
                     'Shape\nSurfaceArea',
                     'Shape\nMajorAxisLength',
                     'T1C\nFirst Order\nRange',
                     'Shape\nMaximum3DDiameter',
                     'T1C\nGLRLM\nGrayLevelNonUniformity'
                     ]

    small_df_plot = df_plot.loc[df_plot['Features'].isin(keep_keys), :]

    fig = plt.figure(figsize=(10, 6))
    ax = fig.add_axes([0.08, 0.4, 0.80, 0.59])

    # Plot
    g = sns.boxplot(x='Features', y='Value', hue='RT', data=small_df_plot, ax=ax, palette="Set1")

    # Rotate feature names
    g.set_xticklabels(g.get_xticklabels(), rotation=60, fontsize=10)
    g.set_xlabel('Features', fontsize=10)
    g.set_ylabel('Normalized Value')

    # Fix legend
    plt.legend(bbox_to_anchor=(1.01, 1), loc=2, borderaxespad=0., title='Recurrence')

    # Save figure
    fig.savefig(os.path.join(spath, 'ttest_feature_box_hand_select_%s.svg' % area), format='svg')
    # fig.savefig(os.path.join(spath, 'ttest_feature_box_hand_select_%s.png' % area), format='png')

    # Close log file
    f.close()


def change_with_RT_rec(df, spath, area='tumor'):
    """
    Computes the difference between pre/post RT and generates plots showing these values.
    Args:
        df (pandas dataframe): radiomic features
        spath (str): savepath
        area (str): mask

    Returns:

    """

    sns.set_style('whitegrid')

    feature_seletion = 'ttest'

    if 'ttest' in feature_seletion:
        # Set up save paths
        spath = os.path.join(spath, 'Analysis', 'ttests')
        if not os.path.exists(spath):
            os.makedirs(spath)

        # Compute p=values
        # Get indicies of pre/post RT
        inds_1 = df['Group'].str.contains('Pre')
        inds_2 = df['Group'].str.contains('Post')
        keys = df.keys()
        non_data_labels = ['Group', 'recurrence', 'rec_days', 'ID']

        # Create a dictionary (later a Dataframe) for evaluation and plotting with Seaborn
        dfn = {'Features': [], 'P-Value': [], 'OrigFeatures': [], 'Contrast': []}

        # Populate dictionary, df, with features and values
        for key in keys:

            if key not in non_data_labels:

                # Get feature values
                tmp = df[key].to_numpy()

                p1 = tmp[inds_1]
                p2 = tmp[inds_2]

                results = stats.ttest_rel(p1, p2)

                if not np.isnan(results.pvalue):

                    # Remove "original" from the plots
                    nkey = key.replace('original_', '')

                    # Capitalize method
                    ind1 = nkey.find('_')
                    method = nkey[:ind1]
                    if method == 'firstorder':
                        method = 'First Order'
                    elif method == 'shape':
                        method = method.capitalize()
                    else:
                        method = method.upper()

                    # Get modality
                    ind2 = nkey.rfind('_')
                    modality = nkey[ind2+1:]
                    if method != 'Shape':
                        nkey = '%s\n%s\n%s' % (modality, method, nkey[ind1+1:ind2])
                    else:
                        nkey = '%s\n%s' % (method, nkey[ind1+1:ind2])

                    # Add contrast
                    if '_T1C' in key:
                        con = 'T1C'

                    elif 'T1' in key:
                        con = 'T1'

                    elif 'T2' in key:
                        con = 'T2'

                    else:
                        con = None

                    if not nkey in dfn['Features']:
                        # Append features, values, recurrence, and IDs to the dictionary
                        dfn['Features'].append(nkey)
                        dfn['OrigFeatures'].append(key)
                        dfn['P-Value'].append(results.pvalue)
                        dfn['Contrast'].append(con)

        # Convert dictionary to Dataframe
        dfn = pd.DataFrame.from_dict(dfn)

        from statsmodels.sandbox.stats.multicomp import multipletests
        mod = multipletests(dfn['P-Value'], alpha=0.05)

        # Remove all but the smallest p-values
        dfn['P-Value'] = mod[1]
        dfn = dfn[mod[0]]
        dfn = dfn.sort_values('P-Value', ascending=True)  # Sort by P-Value, ascending
        dfn = dfn.reset_index(drop=True)

        # Get a list of good features
        features = dfn['OrigFeatures'].to_list()
        # features = features
        n_features = len(features)


    # Set up save path
    spath_tmp = os.path.join(spath, 'Diff_{}'.format(area))
    if not os.path.exists(spath_tmp):
        os.makedirs(spath_tmp)

    # Get all unique ids
    uniq_ids = df['ID'].unique()

    # Features to plot
    features_to_plot = ['original_shape_SurfaceArea_T1',
                        'original_firstorder_TotalEnergy_T1C',
                        'original_gldm_DependenceEntropy_T2',
                        'original_glrlm_LongRunLowGrayLevelemphasis_T2',
                        'original_glrlm_GrayLevelNonUniformity_T1C',
                        'original_glrlm_RunVariance_T1'
                        ]

    # Compute difference (f_post - f_pre) / f_post
    diff_df_plot = {'Feature': [], 'Value': [], 'recurrence': [], 'RT': []}   # Plotting df
    diff_df = {'recurrence': []}                                    # Processing df
    for id in uniq_ids:

        # Get animal data
        animal = df.loc[df['ID'] == id]

        # Get recurrence
        rec = int(animal.loc[['Pre' in a for a in animal['Group']]]['recurrence'])

        # Append to processing dataframe
        diff_df['recurrence'].append(rec)

        # Perform
        for key in features_to_plot:

            # If the key is a data label
            if key not in non_data_labels and key in features[:n_features]:
                # Get pre/post values
                post_val = float(animal.loc[['Post' in a for a in animal['Group']]][key])
                pre_val = float(animal.loc[['Pre' in a for a in animal['Group']]][key])

                # Perform difference
                val = 100 * (post_val - pre_val) / (post_val + 1e-6)


                diff_df_plot['Value'].append(pre_val)
                diff_df_plot['RT'].append('Pre')

                diff_df_plot['Value'].append(post_val)
                diff_df_plot['RT'].append('Post')

                # Clean key for plotting
                key = clean_key(key)

                for _ in range(2):

                    diff_df_plot['Feature'].append(key)

                    # Keep track of recurrence
                    if rec == 0:
                        diff_df_plot['recurrence'].append('No')
                    else:
                        diff_df_plot['recurrence'].append('Yes')

                    if key in diff_df.keys():
                        diff_df[key].append(val)
                    else:
                        diff_df[key] = [val]

    diff_df_plot = pd.DataFrame.from_dict(diff_df_plot)
    # diff_df = pd.DataFrame.from_dict(diff_df)

    # Normalize data
    features_to_plot = [clean_key(i) for i in features_to_plot]
    for key in features_to_plot:

        inds = diff_df_plot['Feature'] == key
        vals = diff_df_plot.loc[inds]['Value']
        mn = vals.min()
        mx = vals.max()

        vals = (vals - mn) / (mx - mn + 1e-6)

        diff_df_plot['Value'][inds] = vals[inds]

    # diff_df_plot = diff_df_plot[keep_keys]

    # Set up figure
    # fig = plt.figure(figsize=(60, 10))
    # ax = fig.add_axes([0.05, 0.35, 0.9, 0.60])
    #
    # # Plot all
    # g = sns.boxplot(x='Feature', y='Value', hue='recurrence', data=diff_df_plot, ax=ax, palette="Set1")
    #
    # # Rotate feature names
    # g.set_xticklabels(g.get_xticklabels(), rotation=35, fontsize=10)
    # g.set_yticklabels(g.get_yticklabels(), fontsize=10)
    # g.set_xlabel('Features', fontsize=10)
    # ax.set_ylabel('Change from Pre to Post RT [%]')
    # plt.ylim([-100, 100])
    #
    # # Fix legend
    # plt.legend(bbox_to_anchor=(1.01, 1), loc=2, borderaxespad=0., title='Recurrence')
    #
    # plt.savefig(os.path.join(spath_tmp, 'diff_all_%s.svg' % area))
    #
    # plt.close(fig)

    # Set up figure
    fig = plt.figure(figsize=(9, 6))
    ax1 = plt.subplot(212)
    ax2 = plt.subplot(211, sharex=ax1)

    df_tmp = diff_df_plot.loc[diff_df_plot['RT'] == 'Post']
    g1 = sns.boxplot(x='Feature', y='Value', hue='recurrence', data=df_tmp, ax=ax1, palette="Set1")
    df_tmp = diff_df_plot.loc[diff_df_plot['RT'] == 'Pre']
    g2 = sns.boxplot(x='Feature', y='Value', hue='recurrence', data=df_tmp, ax=ax2, palette="Set1")
    ax1.get_legend().remove()
    ax2.get_legend().remove()

    # # Rotate feature names
    g1.set_xticklabels(ax1.get_xticklabels(), rotation=40, fontsize=10)
    # ax1.set_yticklabels(ax1.get_yticklabels(), fontsize=10)
    g1.set_xlabel('Features', fontsize=10)
    g1.set_ylabel('Normalized Value', fontsize=10)
    g2.set_xlabel('Features', fontsize=10, visible=False)
    g2.set_ylabel('Normalized Value', fontsize=10)
    g2.set_xticklabels(g1.get_xticklabels(), rotation=35, fontsize=10, visible=False)
    # ax2.set_yticklabels(ax1.get_yticklabels(), fontsize=10)
    ax1.set_ylim([-0.05, 1.05])
    ax2.set_ylim([-0.05, 1.05])

    # Fix legend
    plt.legend(bbox_to_anchor=(1.01, 1), loc=2, borderaxespad=0., title='Recurrence')
    plt.tight_layout()

    plt.savefig(os.path.join(spath_tmp, 'RT_rec_boxplots_%s.svg' % area))

    plt.close(fig)

    # Attempt classification

    # Get numeric data
    # keys = diff_df.keys().tolist()
    # lab_keys = 'recurrence'
    # data_keys = keys
    # data_keys.remove(lab_keys)
    #
    # # Separate data and labels
    # rad = diff_df[data_keys].copy()
    # rad = scale_features(rad)
    # label = diff_df[lab_keys]
    # # ids = diff_df['ID']
    #
    # # SVM classifier
    # svm_classifier_loc(rad, label, spath_tmp)
    #
    # # NN classifier
    # nn_classifier_loc(rad, label, spath_tmp)


def load_data(base_path):

    spath = os.path.join(base_path, 'Analysis')

    snames = [os.path.join(spath, 'tumor_df.csv'),
              os.path.join(spath, 'edge_df.csv'),
              os.path.join(spath, 'bed_df.csv')]

    if not all([os.path.exists(i) for i in snames]):

        # Paths
        summary_file = os.path.join(base_path, 'Summary.xlsx')
        recurrence_file = os.path.join(base_path, 'recurrence.xlsx')
        log_file = os.path.join(base_path, 'processing_log.json')

        recurrence_threshold = 20

        # Animals to exclude from analysis until segmentations are fixed
        exclude = ['K520719', 'K520918', 'K521092', 'K521104', 'K520762']

        # Load summary data
        summary_df, log = load_study_logs(summary_file, log_file)
        df_rec = load_recurrence_log(recurrence_file, recurrence_threshold)

        # Sort data
        df_rec = sort_study_data(df_rec, summary_df, exclude)
        keep = df_rec['animalID'].tolist()

        # Load pre
        radiomics_paths = os.path.join(base_path, 'Radiomics_control_preRT.txt')
        df_pre_tumor_cnt = load_radiomics(radiomics_paths, exclude, region='tumor', keep=keep, group='PreCnt')
        df_pre_edge_cnt = load_radiomics(radiomics_paths, exclude, region='edge', keep=keep, group='PreCnt')
        df_pre_bed_cnt = load_radiomics(radiomics_paths, exclude, region='bed', keep=keep, group='PreCnt')

        radiomics_paths = os.path.join(base_path, 'Radiomics_PD1_preRT.txt')
        df_pre_tumor_pd1 = load_radiomics(radiomics_paths, exclude, region='tumor', keep=keep, group='PrePD1')
        df_pre_edge_pd1 = load_radiomics(radiomics_paths, exclude, region='edge', keep=keep, group='PrePD1')
        df_pre_bed_pd1 = load_radiomics(radiomics_paths, exclude, region='bed', keep=keep, group='PrePD1')

        # Load post
        radiomics_paths = os.path.join(base_path, 'Radiomics_control_postRT.txt')
        df_post_tumor_cnt = load_radiomics(radiomics_paths, exclude, region='tumor', keep=keep, group='PostCnt')
        df_post_edge_cnt = load_radiomics(radiomics_paths, exclude, region='edge', keep=keep, group='PostCnt')
        df_post_bed_cnt = load_radiomics(radiomics_paths, exclude, region='bed', keep=keep, group='PostCnt')

        radiomics_paths = os.path.join(base_path, 'Radiomics_PD1_postRT.txt')
        df_post_tumor_pd1 = load_radiomics(radiomics_paths, exclude, region='tumor', keep=keep, group='PostPD1')
        df_post_edge_pd1 = load_radiomics(radiomics_paths, exclude, region='edge', keep=keep, group='PostPD1')
        df_post_bed_pd1 = load_radiomics(radiomics_paths, exclude, region='bed', keep=keep, group='PostPD1')

        # Concatenate post for RT analysis
        df_tumor = pd.concat((df_pre_tumor_cnt, df_pre_tumor_pd1, df_post_tumor_cnt, df_post_tumor_pd1), ignore_index=True)
        df_edge = pd.concat((df_pre_edge_cnt, df_pre_edge_pd1, df_post_edge_cnt, df_post_edge_pd1), ignore_index=True)
        df_bed = pd.concat((df_pre_bed_cnt, df_pre_bed_pd1, df_post_bed_cnt, df_post_bed_pd1), ignore_index=True)

        # Update recurrence data to account for incomplete animal sets
        uniq = df_tumor['ID'].unique()
        uniq_rec = df_rec['animalID'].unique()
        diff = list(set(uniq_rec).difference(uniq))
        for un in diff:
            df_rec = df_rec[df_rec['animalID'] != un]

        # Check that the sets contains the same animals
        uniq = df_tumor['ID'].unique()
        uniq_rec = df_rec['animalID'].unique()
        diff = list(set(uniq_rec).difference(uniq))
        print('Len of recurrence data: %d' % len(uniq_rec))
        print('Len of image data:      %d' % len(uniq))
        print('Difference: ', list(set(uniq_rec).difference(uniq)))

        # Combine recurrence and radiomics dataframes
        df_tumor = append_rec(df_rec, df_tumor)
        df_edge = append_rec(df_rec, df_edge)
        df_bed = append_rec(df_rec, df_bed)

        # Save dataframes
        df_tumor.to_csv(snames[0])
        df_edge.to_csv(snames[1])
        df_bed.to_csv(snames[2])

    else:

        df_tumor = pd.DataFrame.from_csv(snames[0])
        df_edge = pd.DataFrame.from_csv(snames[1])
        df_bed = pd.DataFrame.from_csv(snames[2])

    return df_tumor, df_edge, df_bed


class PPT:
    def __init__(self, ppt_file):

        self.ppt_file = ppt_file
        self.title_slide = 0
        self.subtitle_slide = 2
        self.title_and_content = 5
        if not os.path.exists(ppt_file):
            prs = pptx.Presentation()

        else:
            prs = pptx.Presentation(ppt_file)

        title_slide_layout = prs.slide_layouts[self.title_slide]
        slide = prs.slides.add_slide(title_slide_layout)

        title = slide.shapes.title
        title.text = 'Radiomic classifications'

        subtitle = slide.placeholders[1]
        subtitle.text = 'Created %s' % datetime.strftime(datetime.now(), '%B %d, %Y')

        self.prs = prs

        # Output file
        self.tmp_file = 'tmp_class.png'


    def add_slides(self, folders, nets=['NN', 'SVM'], im_files=['NN_NoPCA_roc.png', 'SVM_NoPCA_roc.png'], num_features=None):
        """

        Args:
            folders (list of str): path to 1. tumor, 2. edge, and 3. bed

        Returns:

        """

        # Image file names
        # im_files = ['NN_NoPCA_roc.png', 'SVM_NoPCA_roc.png']

        # Get general information
        # Grab folder name for getting group information
        folder_name = os.path.split(folders[0])[1]
        params = folder_name.split(sep='_')

        # Capitalize first letters
        params = [i.lower().capitalize() for i in params]

        # Get group name
        group = params[1]
        mrmr_group = 'mRMR features: %s RT, %s' % (tuple(params[4:]))

        z = 0
        area = []
        files = dict(zip(nets, [[] for _ in nets]))
        for folder in folders:
            # Grab folder name for getting group information
            folder_name = os.path.split(folder)[1]
            params = folder_name.split(sep='_')

            # Capitalize first letters
            params = [i.lower().capitalize() for i in params]

            # Get group name
            area.append(params[2])

            # Get image names
            for i, key in enumerate(nets):
                files[key].append(os.path.join(folder, im_files[i]))

        # Create a single output images
        from PIL import Image, ImageDraw
        offset = 60
        for z in range(len(files['NN'])):
            if z == 0:
                im = Image.open(files['NN'][z])
                height = im.height
                width = im.width

                full_im = Image.new(mode='RGBA', size=(3 * width + offset, len(im_files) * height), color=(255, 255, 255, 255))#(255,255,255,0))

            # Load images
            for i, key in enumerate(nets):
                tmp = Image.open(files[key][z])

                # Compute image coordinates
                bbox = (z * width + offset, i * height)

                # Paste into the larger image
                full_im.paste(tmp, box=bbox)

        # Add text - horizontal
        font = ImageFont.truetype("/usr/share/fonts/truetype/freefont/FreeMono.ttf", 60)
        overlay = Image.new('RGBA', full_im.size, (0,0,0,0))
        draw = ImageDraw.Draw(overlay)
        for z  in range(3):
            bbox1 = (z * width + offset, 45)
            w, h = draw.textsize(area[z], font=font)
            draw.text((bbox1[0] + (width - w) // 2 + 10, bbox1[1]), text=area[z], fill='black', font=font)

        # Add text - vertical
        overlay = overlay.rotate(-90, expand=1)
        draw = ImageDraw.Draw(overlay)
        for z in range(len(nets)):
            bbox1 = (z * height, 10)
            w, h = draw.textsize(area[z], font=font)
            draw.text((bbox1[0] + (height - w)//2, bbox1[1]), text=nets[len(nets)-z-1], fill='black', font=font)
        overlay = overlay.rotate(90, expand=1)

        # Create overlayed image
        out = Image.alpha_composite(full_im, overlay)

        # Save the temporary large image
        out.save(self.tmp_file)

        # Create slide title
        if not num_features:
            section_title = '%s RT. %s' % (group, mrmr_group)
        else:
            section_title = '%s RT. %s, x%d' % (group, mrmr_group, num_features)


        # Update PowerPoint
        content_slide = self.prs.slide_layouts[self.title_and_content]
        slide = self.prs.slides.add_slide(content_slide)

        title = slide.shapes.title
        title.text = section_title

        # Insert image to PPT
        left = pptx.util.Inches(0.1)
        top = pptx.util.Inches(1.6)
        width = pptx.util.Inches(9.8)
        height = pptx.util.Inches(5.7)
        pic = slide.shapes.add_picture(self.tmp_file, left, top, width, height)

    def save(self):

        # Save PPT
        self.prs.save(self.ppt_file)

        # Remove temp image file
        os.remove(self.tmp_file)


def recur_survival(df_tumor, base_path):

    spath = os.path.join(base_path, 'Analysis', 'Survival')
    if not os.path.exists(spath):
        os.makedirs(spath)

    # Set up fitter
    kmf = KaplanMeierFitter()
    df = df_tumor.loc[['Pre' in i for i in df_tumor['Group']]]

    fig, ax = plt.subplots()

    # Non recurrent
    T = df.loc[~df['recurrence']]['rec_days']
    E = T > 180  #np.ones(T.shape) #df_tumor['recurrence']

    kmf.fit(durations=T, event_observed=E, label='No recurrence')
    kmf.plot(ax=ax)

    # Recurrent
    T = df.loc[df['recurrence']]['rec_days']
    E = T > 180  #np.ones(T.shape) #df_tumor['recurrence']

    kmf.fit(durations=T, event_observed=E, label='Recurrence')
    kmf.plot(ax=ax)

    plt.savefig(os.path.join(spath, 'survival.svg'))

def main():

    tstart = time()
    base_path = '/media/matt/Seagate Expansion Drive/b7TData_19/b7TData/Results'

    df_tumor, df_edge, df_bed = load_data(base_path)

    # Survival curves
    # recur_survival(df_tumor, base_path)

    """ 
    A look at differences
    """
    change_with_RT_rec(df_tumor, spath=base_path, area='tumor')
    change_with_RT_rec(df_bed, spath=base_path, area='bed')
    change_with_RT_rec(df_edge, spath=base_path, area='edge')

    # Paired t-test for RT
    paired_ttests_loc(df_tumor, base_path=base_path, area='tumor')
    paired_ttests_loc(df_bed, base_path=base_path, area='bed')
    paired_ttests_loc(df_edge, base_path=base_path, area='edge')


    # mRMR - Test all combinations
    groups = ['Pre', 'Post']
    mrmr_files = ['mRMR_results_%s_tumor.txt', 'mRMR_results_%s_bed.txt',
                  'mRMR_results_%s_edge.txt']

    # Set up PPT
    ppt_file = os.path.join(base_path, 'Analysis', 'classifications3.pptx')
    ppt = PPT(ppt_file)

    spaths = [''] * 3
    num_features = range(10, 115)
    mrmr_file = 'mRMR_results_%s_edge.txt'

    gropus = ['Pre']
    mrmr_files = ['mRMR_results_%s_edge.txt']
    for group in groups:
        for mrmr_file in mrmr_files:
            nf = 108
            mrmr_file = mrmr_file % group
            spaths[0] = mRMR(df_tumor, group=group, base_path=base_path, area='tumor', mrmr_file=mrmr_file,
                             num_features=nf, htmaps=True)
            spaths[1] = mRMR(df_edge, group=group, base_path=base_path, area='edge', mrmr_file=mrmr_file,
                             num_features=nf, htmaps=True)
            spaths[2] = mRMR(df_bed, group=group, base_path=base_path, area='bed', mrmr_file=mrmr_file,
                             num_features=nf, htmaps=True)

            # Add classifications to a PPT
            ppt.add_slides(folders=spaths, num_features=nf)

            ppt.save()


    # Test Post RT
    ppt_file = os.path.join(base_path, 'Analysis', 'classifications4.pptx')
    ppt = PPT(ppt_file)

    spaths = [''] * 3
    num_features = range(10, 115)
    # mrmr_file = 'mRMR_results_%s_edge.txt'

    groups = ['Post']
    htmaps = False
    # mrmr_files = ['mRMR_results_%s_edge.txt']
    for group in groups:
        for mrmr_file in mrmr_files:
            nf = 108
            mrmr_file = mrmr_file % group
            spaths[0] = mRMR(df_tumor, group=group, base_path=base_path, area='tumor', mrmr_file=mrmr_file,
                             num_features=nf, htmaps=htmaps)
            spaths[1] = mRMR(df_edge, group=group, base_path=base_path, area='edge', mrmr_file=mrmr_file,
                             num_features=nf, htmaps=htmaps)
            spaths[2] = mRMR(df_bed, group=group, base_path=base_path, area='bed', mrmr_file=mrmr_file,
                             num_features=nf, htmaps=htmaps)

            # Add classifications to a PPT
            ppt.add_slides(folders=spaths, num_features=nf)

            ppt.save()

    print('\tTotal time (HH:MM:SS): %s\n\n' % (str(dt.timedelta(seconds=round(time() - tstart)))))



if __name__ == '__main__':

    main()

    folders = ['/media/matt/Seagate Expansion Drive/b7TData_19/b7TData/Results/Analysis/mrmr_Pre_tumor_feat_Pre_tumor',
               '/media/matt/Seagate Expansion Drive/b7TData_19/b7TData/Results/Analysis/mrmr_Pre_edge_feat_Pre_tumor',
               '/media/matt/Seagate Expansion Drive/b7TData_19/b7TData/Results/Analysis/mrmr_Pre_bed_feat_Pre_tumor']
